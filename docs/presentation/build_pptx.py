#!/usr/bin/env python3
"""Build the final-presentation .pptx from the same content as index.html.

Design intent: dark-navy academic deck. No rounded corners. No em dashes.
IBM Plex Serif on display headings, Inter for body, JetBrains Mono for labels.
8 slides @ 1920x1080.
"""
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree
import os

HERE = os.path.dirname(os.path.abspath(__file__))
DOCS = os.path.dirname(HERE)
ARCH_PNG = os.path.join(DOCS, 'diagrams', 'Architecture.png')
OUT = os.path.join(HERE, 'Presentation.pptx')

W = Emu(1920 * 9525)   # 1920 px @ 96dpi -> EMU
H = Emu(1080 * 9525)

# ---- palette (mirrors CSS variables) -----------------------------------
BG       = RGBColor(0x07, 0x09, 0x1a)   # --bg-0
INK_0    = RGBColor(0xf1, 0xf5, 0xf9)
INK_1    = RGBColor(0xcb, 0xd5, 0xe1)
INK_2    = RGBColor(0x94, 0xa3, 0xb8)
INK_3    = RGBColor(0x64, 0x74, 0x8b)
RULE     = RGBColor(0x2b, 0x32, 0x47)   # rgba(148,163,184,0.20) flattened on bg
ACCENT   = RGBColor(0x38, 0xbd, 0xf8)
ACCENT_2 = RGBColor(0xc0, 0x84, 0xfc)

SERIF = 'IBM Plex Serif'
SANS  = 'Inter'
MONO  = 'JetBrains Mono'

PX = 9525  # EMU per CSS px at 96 dpi


def px(n):
    return Emu(int(n * PX))


# ---- helpers ----------------------------------------------------------
def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=Pt(0.75)):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = line_w
    # remove default rounded corners (none on rectangles, but ensure no shadow)
    return shp


def add_line(slide, x, y, w, color=RULE, weight=0.75):
    """A thin horizontal rule of width w starting at (x, y)."""
    shp = slide.shapes.add_connector(1, x, y, x + w, y)  # 1 = STRAIGHT
    shp.line.color.rgb = color
    shp.line.width = Pt(weight)
    return shp


def add_text(slide, x, y, w, h, runs, *, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT,
             margin=0, line_spacing=None):
    """runs: list of (text, kwargs) where kwargs may include
       font, size_pt, bold, color, italic. A run text of "\n" forces a new paragraph."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(margin)
    tf.margin_top = tf.margin_bottom = Emu(margin)
    tf.vertical_anchor = anchor

    # First paragraph already exists; we add to it then add new ones for "\n"
    paragraphs = [tf.paragraphs[0]]
    paragraphs[0].alignment = align

    def cur():
        return paragraphs[-1]

    first = True
    for text, opts in runs:
        if text == '\n':
            p = tf.add_paragraph()
            p.alignment = align
            if line_spacing is not None:
                p.line_spacing = line_spacing
            paragraphs.append(p)
            first = True
            continue
        p = cur()
        if first and line_spacing is not None:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = text
        f = run.font
        f.name = opts.get('font', SANS)
        if 'size_pt' in opts:
            f.size = Pt(opts['size_pt'])
        if 'bold' in opts:
            f.bold = bool(opts['bold'])
        if 'italic' in opts:
            f.italic = bool(opts['italic'])
        if 'color' in opts and opts['color'] is not None:
            f.color.rgb = opts['color']
        first = False
    return tb


def add_para_text(slide, x, y, w, h, text, *, font=SANS, size_pt=18, color=INK_1,
                  bold=False, line_spacing=1.4, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT):
    return add_text(slide, x, y, w, h, [(text, dict(font=font, size_pt=size_pt, color=color, bold=bold))],
                    anchor=anchor, align=align, line_spacing=line_spacing)


# ---- universal slide chrome -------------------------------------------
def slide_chrome(slide, top_left_label, page_label, footer_left, footer_right):
    # Top rule + meta
    add_text(slide, px(56), px(44), px(1920 - 112), px(28), [
        (top_left_label, dict(font=MONO, size_pt=10, color=INK_2, bold=True)),
    ])
    add_text(slide, px(56), px(44), px(1920 - 112), px(28), [
        (page_label, dict(font=MONO, size_pt=10, color=INK_3, bold=True)),
    ], align=PP_ALIGN.RIGHT)
    add_line(slide, px(56), px(82), px(1920 - 112))

    # Bottom rule + footer
    add_line(slide, px(56), px(1080 - 56), px(1920 - 112))
    add_text(slide, px(56), px(1080 - 50), px(1920 - 112), px(20), [
        (footer_left, dict(font=MONO, size_pt=9, color=INK_3, bold=True)),
    ])
    add_text(slide, px(56), px(1080 - 50), px(1920 - 112), px(20), [
        (footer_right, dict(font=MONO, size_pt=9, color=INK_3, bold=True)),
    ], align=PP_ALIGN.RIGHT)


FOOT_L = 'NYU CLOUD COMPUTING  •  SPRING 2026'
FOOT_R = 'MISHRA  •  ZHENG  •  SANKPAL  •  HUANG'


# ----- build ------------------------------------------------------------
prs = Presentation()
prs.slide_width = W
prs.slide_height = H
blank = prs.slide_layouts[6]


# ===== Slide 1 — Title =====
s = prs.slides.add_slide(blank)
set_bg(s, BG)

# Tiny institution rule
add_text(s, px(96), px(160), px(1500), px(28), [
    ('NEW YORK UNIVERSITY  •  DEPARTMENT OF COMPUTER SCIENCE  •  CLOUD COMPUTING  •  SPRING 2026',
     dict(font=MONO, size_pt=10, color=INK_2, bold=True)),
])
add_line(s, px(96), px(200), px(1100))

# Display title
add_text(s, px(96), px(240), px(1700), px(360), [
    ('Research Paper\nSummarizer.', dict(font=SERIF, size_pt=92, color=INK_0, bold=False)),
], line_spacing=0.96)

# Accent rule
shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(96), px(620), px(200), px(3))
shp.fill.solid(); shp.fill.fore_color.rgb = ACCENT
shp.line.fill.background()

# Subtitle
add_para_text(s, px(96), px(660), px(1500), px(220),
              ('A serverless, event-driven AWS platform that searches arXiv and Semantic Scholar, '
               'fetches full-text PDFs, and produces structured five-section summaries in roughly thirty seconds. '
               'Built end to end on managed services with zero idle cost.'),
              font=SANS, size_pt=20, color=INK_1, line_spacing=1.42)

# Footer block (top rule + authors / url)
add_line(s, px(96), px(940), px(1920 - 192))
add_text(s, px(96), px(960), px(900), px(80), [
    ('Pranjal Mishra  •  Yang Zheng  •  Shreyas Sankpal  •  Kerry Huang',
     dict(font=SANS, size_pt=15, color=INK_1, bold=True)),
    ('\n', {}),
    ('FINAL PROJECT', dict(font=MONO, size_pt=9, color=INK_3, bold=True)),
], line_spacing=1.4)
add_text(s, px(900), px(960), px(900), px(40), [
    ('d24irdkbe9jj2b.cloudfront.net', dict(font=MONO, size_pt=12, color=ACCENT, bold=True)),
], align=PP_ALIGN.RIGHT)


# ===== helper to start a normal slide =====
def begin(top_label, page):
    s = prs.slides.add_slide(blank)
    set_bg(s, BG)
    slide_chrome(s, top_label, page, FOOT_L, FOOT_R)
    return s


def section_heading(slide, text, *, y=130, h=160):
    add_text(slide, px(96), px(y), px(1920 - 192), px(h), [
        (text, dict(font=SERIF, size_pt=40, color=INK_0)),
    ], line_spacing=1.1)


def lede(slide, text, *, y=240, w=1500):
    add_para_text(slide, px(96), px(y), px(w), px(180), text, font=SANS, size_pt=16, color=INK_1, line_spacing=1.5)


# ===== Slide 2 — System overview =====
s = begin('I.  SYSTEM OVERVIEW', '01 / 08')

# Display
add_text(s, px(96), px(130), px(1920 - 192), px(160), [
    ('Triage is the bottleneck.', dict(font=SERIF, size_pt=82, color=INK_0)),
], line_spacing=1.0)
lede(s, ('Search engines surface relevant titles. Citation graphs reveal influence. '
         'General-purpose chat assistants will summarize a paper if the user pastes it. '
         'None of these, alone, is a workflow.'),
     y=280, w=1500)

# Stat row
add_line(s, px(96), px(420), px(1920 - 192))
add_text(s, px(96), px(450), px(420), px(160), [
    ('2.4M+', dict(font=SERIF, size_pt=80, color=ACCENT)),
], line_spacing=1.0)
add_para_text(s, px(540), px(470), px(1240), px(160),
              ('preprints indexed by arXiv alone, doubling roughly every decade. '
               'The platform delivers three composable capabilities on one shared backbone of chunk vectors persisted in Amazon DynamoDB.'),
              font=SANS, size_pt=17, color=INK_1, line_spacing=1.5)

# Three pillars (rectangle borders only)
add_line(s, px(96), px(660), px(1920 - 192))
col_w = (1920 - 192) // 3
pillars = [
    ('01.  SEARCH', 'Multi-source ranked search.',
     'arXiv and Semantic Scholar invoked in parallel from a single Lambda; results deduplicated by arxivId, DOI, and normalized title in one to two seconds end to end.'),
    ('02.  SUMMARIZE', 'Asynchronous five-section summary.',
     'Objectives, methodology, results, limitations, contributions. SQS-triggered Step Functions map-reduce. Done in roughly thirty seconds for a typical paper.'),
    ('03.  GROUND', 'Retrieval-augmented features.',
     'Talk-to-PDF chat with cited chunks, related-papers ranking by cosine similarity, and a typed knowledge graph; all from the embeddings the pipeline already wrote.'),
]
for i, (num, head, body) in enumerate(pillars):
    x0 = px(96 + i * col_w)
    if i > 0:
        # vertical divider
        sep = s.shapes.add_connector(1, x0, px(680), x0, px(940))
        sep.line.color.rgb = RULE; sep.line.width = Pt(0.75)
    add_text(s, x0 + px(20), px(700), px(col_w - 40), px(28), [
        (num, dict(font=MONO, size_pt=10, color=ACCENT, bold=True)),
    ])
    add_text(s, x0 + px(20), px(730), px(col_w - 40), px(50), [
        (head, dict(font=SERIF, size_pt=20, color=INK_0)),
    ], line_spacing=1.18)
    add_para_text(s, x0 + px(20), px(800), px(col_w - 40), px(160), body,
                  font=SANS, size_pt=13, color=INK_1, line_spacing=1.55)
add_line(s, px(96), px(940), px(1920 - 192))


# ===== Slide 3 — Architecture =====
s = begin('II.  ARCHITECTURE', '02 / 08')
section_heading(s, 'Fifteen single-purpose Lambdas, six CDK stacks, three managed data stores.', y=130)
lede(s, ('Inter-service communication restricted to API Gateway, Amazon SQS, AWS Step Functions, '
         'and Amazon DynamoDB. No Lambda directly invokes another Lambda.'),
     y=280, w=1700)

# Architecture canvas — sized exactly to the image aspect ratio so there is
# no empty white space on the sides. Image AR = 3157/1811 = 1.743.
img_w_px, img_h_px = 3157, 1811
img_ar = img_w_px / img_h_px
pad_in_px = 28
target_h_px = 580                                    # fits between lede and footer
target_inner_h_px = target_h_px - 2 * pad_in_px
target_inner_w_px = int(target_inner_h_px * img_ar)
canvas_w_px = target_inner_w_px + 2 * pad_in_px
canvas_x_px = (1920 - canvas_w_px) // 2
canvas_y_px = 400
canvas_x, canvas_y = px(canvas_x_px), px(canvas_y_px)
canvas_w, canvas_h = px(canvas_w_px), px(target_h_px)
add_rect(s, canvas_x, canvas_y, canvas_w, canvas_h, fill=RGBColor(0xf8, 0xfa, 0xfc), line=RULE)
if os.path.exists(ARCH_PNG):
    s.shapes.add_picture(ARCH_PNG,
                         canvas_x + px(pad_in_px), canvas_y + px(pad_in_px),
                         width=px(target_inner_w_px), height=px(target_inner_h_px))


# ===== Slide 4 — Pipeline =====
s = begin('III.  ASYNCHRONOUS DATA PIPELINE', '03 / 08')
section_heading(s, 'Why a state machine, not one large Lambda.', y=130)
lede(s, ('Map-reduce summarization is the hot path. The five-stage AWS Step Functions execution sidesteps the '
         'Lambda fifteen-minute ceiling, isolates per-chunk failures, and keeps state inputs under the 256 KB limit.'),
     y=230, w=1700)

# 5 stages, equal columns
add_line(s, px(96), px(370), px(1920 - 192))
stage_w = (1920 - 192) // 5
stages = [
    ('Stage 1', 'FetchPDF', 'S3-cached', False),
    ('Stage 2', 'ExtractText', 'pdf-parse', False),
    ('Stage 3', 'Chunk + Embed', '~1500 tok / chunk, 400 char overlap', False),
    ('Stage 4 · Map', 'MapSummarize', 'parallel, capped at 5', True),
    ('Stage 5', 'ReduceSummary', 'DynamoDB write', False),
]
for i, (lbl, name, sub, is_map) in enumerate(stages):
    x0 = px(96 + i * stage_w)
    if i > 0:
        sep = s.shapes.add_connector(1, x0, px(390), x0, px(560))
        sep.line.color.rgb = RULE; sep.line.width = Pt(0.75)
    add_text(s, x0 + px(20), px(410), px(stage_w - 40), px(28), [
        (lbl.upper(), dict(font=MONO, size_pt=9, color=(ACCENT if is_map else INK_3), bold=True)),
    ])
    add_text(s, x0 + px(20), px(442), px(stage_w - 40), px(48), [
        (name, dict(font=SERIF, size_pt=22, color=(ACCENT if is_map else INK_0))),
    ], line_spacing=1.1)
    add_para_text(s, x0 + px(20), px(488), px(stage_w - 40), px(60), sub,
                  font=SANS, size_pt=11, color=INK_2, line_spacing=1.4)
add_line(s, px(96), px(560), px(1920 - 192))

# Why row
add_line(s, px(96), px(660), px(1920 - 192))
why_w = (1920 - 192) // 3
whys = [
    ('Past the fifteen-minute Lambda ceiling.',
     'A thirty-page paper, twelve parallel chunks, plus reduction can blow past it. Step Functions has no such limit; long-running orchestration becomes free.'),
    ('Per-chunk retry, not per-job.',
     'If one Bedrock call is throttled, that chunk retries with exponential backoff. The other eleven succeed once, and exactly once. No duplicate work, no whole-job redo.'),
    ('S3 keys, not text, in state.',
     'Step Functions caps state inputs at 256 KB. Heavy data (PDFs, extracted text, chunks, embeddings) lives in S3. We thread only S3 keys, never bytes.'),
]
for i, (h, body) in enumerate(whys):
    x0 = px(96 + i * why_w)
    if i > 0:
        sep = s.shapes.add_connector(1, x0, px(680), x0, px(900))
        sep.line.color.rgb = RULE; sep.line.width = Pt(0.75)
    add_text(s, x0 + px(20), px(700), px(why_w - 40), px(50), [
        (h, dict(font=SANS, size_pt=15, color=INK_0, bold=True)),
    ], line_spacing=1.18)
    add_para_text(s, x0 + px(20), px(760), px(why_w - 40), px(160), body,
                  font=SANS, size_pt=12, color=INK_1, line_spacing=1.55)
add_line(s, px(96), px(900), px(1920 - 192))


# ===== Slide 5 — Data layer =====
s = begin('IV.  DATA LAYER · SINGLE-TABLE DESIGN', '04 / 08')
section_heading(s, 'One DynamoDB table. One global secondary index. Every access pattern.', y=130)
lede(s, ('Access patterns are narrow and known in advance. A single-table layout with one GSI keyed on a content '
         'hash handles every read, with single-digit-millisecond latency under on-demand billing.'),
     y=290, w=1700)

# Left: DDB schema in a bordered rect
left_x, left_y, left_w, left_h = px(96), px(420), px(960), px(520)
add_rect(s, left_x, left_y, left_w, left_h, line=RULE)

add_text(s, left_x + px(28), left_y + px(24), left_w - px(56), px(40), [
    ('MainTable', dict(font=SERIF, size_pt=20, color=INK_0)),
])
add_text(s, left_x + px(28), left_y + px(58), left_w - px(56), px(28), [
    ('PK = USER#<userId>.   SK prefix encodes entity type.',
     dict(font=MONO, size_pt=10, color=INK_2)),
])

# Headers
hdr_y = left_y + px(110)
add_text(s, left_x + px(28),  hdr_y, px(280), px(20), [('PK',  dict(font=MONO, size_pt=9, color=INK_3, bold=True))])
add_text(s, left_x + px(330), hdr_y, px(220), px(20), [('SK',  dict(font=MONO, size_pt=9, color=INK_3, bold=True))])
add_text(s, left_x + px(560), hdr_y, px(380), px(20), [('ITEM', dict(font=MONO, size_pt=9, color=INK_3, bold=True))])
add_line(s, left_x + px(28), hdr_y + px(24), left_w - px(56))

rows = [
    ('USER#u123', 'PROFILE#', 'quota, plan', INK_0, ACCENT),
    ('USER#u123', 'JOB#j789', 'summary, sections, meanEmbedding', INK_0, ACCENT),
    ('JOB#j789', 'CHUNK#0', 'text, 1024-d Titan v2 vector', INK_0, ACCENT),
    ('JOB#j789', 'CHUNK#1', 'text, vector', INK_0, ACCENT),
    ('GSI1: PAPER#<sha256>', 'JOB#j789', 'cross-user content-hash dedup', ACCENT_2, ACCENT_2),
]
ry = hdr_y + px(40)
for pk, sk, item, _, color in rows:
    add_text(s, left_x + px(28),  ry, px(290), px(28), [(pk, dict(font=MONO, size_pt=11, color=color))])
    add_text(s, left_x + px(330), ry, px(220), px(28), [(sk, dict(font=MONO, size_pt=11, color=INK_1))])
    add_text(s, left_x + px(560), ry, px(380), px(28), [(item, dict(font=MONO, size_pt=11, color=INK_1))])
    add_line(s, left_x + px(28), ry + px(36), left_w - px(56), color=RGBColor(0x1c, 0x21, 0x32))
    ry = ry + px(60)

# Right: 3 cards
right_x = px(1100); right_w = px(728); card_h = px(164)
notes = [
    ('Content-hash deduplication via GSI1.',
     'A second user submitting a paper a first user already summarized triggers a GSI1 lookup, copies the existing JOB record, and refunds the second user\'s quota.'),
    ('Two S3 buckets, distinct postures.',
     'Frontend bucket: private, served only via CloudFront Origin Access Control. PDFs bucket: 30-day IA transition, 90-day expiration. Storage cost stays flat.'),
    ('Encryption and on-demand billing.',
     'AWS-managed KMS at rest. On-demand DynamoDB billing means storage cost scales with actual usage. Block-public-access enabled on every bucket.'),
]
y = px(420)
for i, (h, p) in enumerate(notes):
    add_rect(s, right_x, y, right_w, card_h, line=RULE)
    add_text(s, right_x + px(22), y + px(18), right_w - px(44), px(32), [
        (h, dict(font=SANS, size_pt=15, color=INK_0, bold=True)),
    ], line_spacing=1.2)
    add_para_text(s, right_x + px(22), y + px(58), right_w - px(44), px(120), p,
                  font=SANS, size_pt=12, color=INK_1, line_spacing=1.5)
    y = y + card_h + px(12)


# ===== Slide 6 — Reliability =====
s = begin('V.  RELIABILITY · OBSERVABILITY · IAM', '05 / 08')
section_heading(s, 'Every cross-service call has an explicit failure path.', y=130)
lede(s, ('Idempotent writes, bounded retries, dead-letter quarantine, alarmed observability, least-privilege IAM. '
         'Six properties the rubric asks for and the system actually exhibits.'),
     y=230, w=1700)

# 3x2 grid of plain bordered cells
grid_x, grid_y = px(96), px(370)
grid_w, grid_h = px(1920 - 192), px(560)
cw = grid_w // 3
ch = grid_h // 2

# Outer rectangle
add_rect(s, grid_x, grid_y, grid_w, grid_h, line=RULE)
# Internal dividers
for i in (1, 2):
    sep = s.shapes.add_connector(1, grid_x + cw * i, grid_y, grid_x + cw * i, grid_y + grid_h)
    sep.line.color.rgb = RULE; sep.line.width = Pt(0.75)
sep = s.shapes.add_connector(1, grid_x, grid_y + ch, grid_x + grid_w, grid_y + ch)
sep.line.color.rgb = RULE; sep.line.width = Pt(0.75)

cells = [
    ('01', 'Idempotent by construction.',
     'Step Functions execution name derived from jobId; an SQS redrive is rejected as duplicate. Every DynamoDB write keys on jobId, so a replay never corrupts state.'),
    ('02', 'Retry where it makes sense, DLQ where it does not.',
     'Step Functions retries each task with exponential backoff. SQS redrive policy moves messages to a dead-letter queue after three failed attempts. No infinite loops.'),
    ('03', 'Failures surface, never vanish.',
     'A MarkFailed state catches any uncaught error from any prior state, captures the cause, updates the JOB record to status "failed". The user sees a clean reason.'),
    ('04', 'Observability across every Lambda.',
     'Four CloudWatch alarms (DLQ depth, SF failure count, pipeline errors, API errors) wired to SNS email. AWS X-Ray distributed tracing. AWS Budgets at $50 / month.'),
    ('05', 'Least-privilege IAM, per function.',
     'searchFn cannot read the PDFs bucket. getSummaryFn cannot write DynamoDB. mapSummarize cannot read Cognito. Each Lambda role is the narrowest it can be.'),
    ('06', 'No Lambda-to-Lambda direct calls.',
     'Inter-service communication is API Gateway, Amazon SQS, AWS Step Functions, or Amazon DynamoDB. Never lambda.invoke. A defect cannot cascade by direct call.'),
]
for i, (num, h, body) in enumerate(cells):
    col = i % 3; row = i // 3
    x0 = grid_x + cw * col
    y0 = grid_y + ch * row
    pad = px(22)
    add_text(s, x0 + cw - pad - px(40), y0 + px(18), px(60), px(20), [
        (num, dict(font=MONO, size_pt=10, color=INK_3, bold=True)),
    ], align=PP_ALIGN.RIGHT)
    add_text(s, x0 + pad, y0 + px(28), cw - 2 * pad, px(50), [
        (h, dict(font=SERIF, size_pt=17, color=INK_0)),
    ], line_spacing=1.18)
    add_para_text(s, x0 + pad, y0 + px(85), cw - 2 * pad, ch - px(110), body,
                  font=SANS, size_pt=12, color=INK_1, line_spacing=1.55)


# ===== Slide 7 — Numbers =====
s = begin('VI.  EVALUATION', '06 / 08')
section_heading(s, 'Measured against real arXiv submissions.', y=130)
lede(s, ('Numbers below are drawn from CloudWatch metrics and DynamoDB on the live deployment. '
         'Not benchmarks. Not promises.'),
     y=230, w=1700)

# 4 columns separated by rules
add_line(s, px(96), px(420), px(1920 - 192))
ncw = (1920 - 192) // 4
nums = [
    ('END-TO-END SUMMARY', '27–50', 's', ACCENT,
     'Submit to "done", chunked, parallel-summarized, reduced. Wall-clock dominated by the slowest chunk plus reduction, not by chunk count.'),
    ('MARGINAL COST', '$0.20', '', INK_0,
     'per paper, dominated by Bedrock tokens. Repeat submissions are free via content-hash deduplication on GSI1.'),
    ('IDLE MONTHLY COST', '$0.00', '', INK_0,
     'Every component except Bedrock fits the AWS free tier at our scale. No traffic, no bill. By design, not by accident.'),
    ('DLQ MESSAGES', '0', '', INK_0,
     'Messages stuck in the dead-letter queue across the deployment. Every failure surfaces with a cause string written to the JOB record.'),
]
for i, (lbl, big, suffix, color, body) in enumerate(nums):
    x0 = px(96 + i * ncw)
    if i > 0:
        sep = s.shapes.add_connector(1, x0, px(440), x0, px(900))
        sep.line.color.rgb = RULE; sep.line.width = Pt(0.75)
    add_text(s, x0 + px(28), px(460), px(ncw - 56), px(28), [
        (lbl, dict(font=MONO, size_pt=10, color=INK_3, bold=True)),
    ])
    add_text(s, x0 + px(28), px(498), px(ncw - 56), px(140), [
        (big, dict(font=SERIF, size_pt=64, color=color)),
        (' ' + suffix if suffix else '', dict(font=SERIF, size_pt=24, color=INK_2)),
    ], line_spacing=1.0)
    add_para_text(s, x0 + px(28), px(680), px(ncw - 56), px(220), body,
                  font=SANS, size_pt=12, color=INK_2, line_spacing=1.55)
add_line(s, px(96), px(900), px(1920 - 192))


# ===== Slide 8 — Team =====
s = begin('VII.  TEAM AND RESOURCES', '07 / 08')
section_heading(s, 'Built by four people, owned end to end.', y=130)

# 2x2 team grid
add_line(s, px(96), px(260), px(1920 - 192))
mw_px = (1920 - 192) // 2
mw = px(mw_px)
mh = px(220)
members = [
    ('Pranjal Mishra', 'ARCHITECTURE, FRONTEND, RAG, OPS',
     'Six CDK stacks, the Next.js frontend, the AWS Step Functions state machine, Amazon Bedrock integration, RAG chat with cited chunks, the knowledge graph extractor and visualization, the CloudWatch dashboard, four alarms, AWS X-Ray tracing, AWS Budgets, the deploy automation.'),
    ('Yang Zheng', 'SEARCH SERVICE',
     'searchFn Lambda, the arXiv API client, the Atom XML parser, parameter validation with structured error handling, and 14 vitest unit tests covering query construction, Atom parsing, and handler edge cases.'),
    ('Shreyas Sankpal', 'PIPELINE ORCHESTRATION',
     'Pipeline orchestration design discussions and chunking-strategy review for the map-reduce summarization stages.'),
    ('Kerry Huang', 'AUTHENTICATION AND API',
     'Authentication and API design discussions for the Amazon Cognito User Pool, the Secure Remote Password login flow, and the API Gateway authorizer surface.'),
]

for i, (name, role, body) in enumerate(members):
    col = i % 2; row = i // 2
    x0 = px(96 + col * mw_px)
    y0 = px(260) + row * mh
    if col == 1:
        sep = s.shapes.add_connector(1, x0, y0 + px(8), x0, y0 + mh - px(8))
        sep.line.color.rgb = RULE; sep.line.width = Pt(0.75)
    pad = px(28)
    add_text(s, x0 + pad, y0 + px(20), mw - 2 * pad, px(38), [
        (name, dict(font=SERIF, size_pt=22, color=INK_0)),
    ], line_spacing=1.0)
    add_text(s, x0 + pad, y0 + px(58), mw - 2 * pad, px(20), [
        (role, dict(font=MONO, size_pt=10, color=ACCENT, bold=True)),
    ])
    add_para_text(s, x0 + pad, y0 + px(90), mw - 2 * pad, mh - px(110), body,
                  font=SANS, size_pt=13, color=INK_1, line_spacing=1.55)
add_line(s, px(96), px(700), px(1920 - 192))

# End bar
end_y = px(740)
add_rect(s, px(96), end_y, px(1920 - 192), px(120), line=RULE)
add_text(s, px(120), end_y + px(38), px(900), px(48), [
    ('Thank you.  ', dict(font=SERIF, size_pt=22, color=INK_0)),
    ('Live demonstration follows.', dict(font=SERIF, size_pt=22, color=ACCENT)),
])
add_text(s, px(1920 - 96 - 720), end_y + px(44), px(696), px(40), [
    ('d24irdkbe9jj2b.cloudfront.net    •    github.com/pmxlr8/research-summarizer',
     dict(font=MONO, size_pt=12, color=ACCENT)),
], align=PP_ALIGN.RIGHT)


prs.save(OUT)
print(f'Wrote {OUT}')
print(f'Slides: {len(prs.slides)}')
